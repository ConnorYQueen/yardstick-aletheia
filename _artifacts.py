"""Local artifact builders for Aletheia: decks, documents, financial models,
and the autonomy dashboard.

Aletheia produces a structured spec (JSON) for an artifact; this module renders
it into a real file on the customer's machine, applying their branding and the
confidentiality markings from brand.json. Everything is built LOCALLY with
pure-Python libraries (python-pptx / fpdf2 / openpyxl) - no system libraries,
no network, nothing leaves the customer's environment. That is what makes it
safe for locked-down and regulated environments.

Builders are dependency-light and optional: a customer only needs the library
for the artifact type they ask for. Missing a library returns a clear
"pip install" message instead of crashing the session. The dashboard builder
needs no library at all -- it is pure stdlib HTML.

brand.json (written by Aletheia after she asks for branding + compliance):
    {
      "company": "Acme Corp",
      "primary": "#1F3A5F", "secondary": "#5B7DB1", "accent": "#E08D3C",
      "font": "Helvetica", "logo": "logo.png",
      "confidentiality_footer": "Confidential - Internal Use Only",
      "classification": "INTERNAL"
    }
"""
from __future__ import annotations

import json
import re
from pathlib import Path

# Neutral professional defaults until the customer provides brand.json.
DEFAULT_BRAND = {
    "company": "[Your Company]",
    "primary": "#1F2A37",
    "secondary": "#4B5563",
    "accent": "#C2703D",
    "font": "Helvetica",
    "logo": "",
    "confidentiality_footer": "",
    "classification": "",
}


def load_brand(pack_dir: Path) -> dict:
    b = dict(DEFAULT_BRAND)
    p = pack_dir / "brand.json"
    if p.exists():
        try:
            b.update({k: v for k, v in json.loads(p.read_text(encoding="utf-8")).items()
                      if v not in (None, "")})
        except (json.JSONDecodeError, OSError):
            pass
    return b


def _logo_path(pack_dir: Path, brand: dict) -> Path | None:
    if not brand.get("logo"):
        return None
    lp = (pack_dir / brand["logo"]).resolve()
    # keep the logo reference inside the pack folder
    try:
        lp.relative_to(pack_dir.resolve())
    except ValueError:
        return None
    return lp if lp.exists() else None


def _hex(s: str, default="000000") -> str:
    s = (s or "").lstrip("#").strip()
    return s if re.fullmatch(r"[0-9a-fA-F]{6}", s) else default


def _rgb_tuple(s: str) -> tuple:
    h = _hex(s)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _slug(title: str, fallback: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "-", (title or "").strip().lower()).strip("-")
    return s or fallback


# ---------------------------------------------------------------------------
# Deck (PPTX)
# ---------------------------------------------------------------------------

def build_deck(title: str, spec: dict, brand: dict, pack_dir: Path, out: Path) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        return "[deck skipped: run `pip install python-pptx` to build decks]"

    primary = RGBColor(*_rgb_tuple(brand["primary"]))
    accent = RGBColor(*_rgb_tuple(brand["accent"]))
    ink = RGBColor(0x1F, 0x29, 0x37)
    font = brand.get("font") or "Helvetica"
    company = brand.get("company", "")
    footer = brand.get("confidentiality_footer", "")
    classif = brand.get("classification", "")
    logo = _logo_path(pack_dir, brand)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    EMU_W = Inches(13.333)

    def text(slide, x, y, w, h, runs, size, color, bold=False, align_right=False):
        from pptx.enum.text import PP_ALIGN
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(runs if isinstance(runs, list) else [runs]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(line)
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.name = font
            p.font.color.rgb = color
            if align_right:
                p.alignment = PP_ALIGN.RIGHT
        return tb

    def chrome(slide):
        # accent bar + footer + classification on every slide
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.32), EMU_W, Inches(0.18))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        foot = " / ".join(p for p in (company, footer) if p)
        if foot:
            text(slide, Inches(0.4), Inches(7.0), Inches(9), Inches(0.35), foot, 9, RGBColor(0x6B, 0x72, 0x80))
        if classif:
            text(slide, Inches(9.5), Inches(0.15), Inches(3.4), Inches(0.3), classif.upper(), 9, accent, bold=True, align_right=True)

    # Cover
    s = prs.slides.add_slide(blank)
    chrome(s)
    if logo:
        try: s.shapes.add_picture(str(logo), Inches(0.5), Inches(0.5), height=Inches(0.7))
        except Exception: pass
    text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(1.6), title, 40, primary, bold=True)
    if spec.get("subtitle"):
        text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(1.0), spec["subtitle"], 18, ink)

    # Content slides
    for sl in spec.get("slides", []):
        s = prs.slides.add_slide(blank)
        chrome(s)
        text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.9), sl.get("heading", ""), 26, primary, bold=True)
        y = Inches(1.7)
        if sl.get("body"):
            text(s, Inches(0.6), y, Inches(12.1), Inches(1.5), sl["body"], 14, ink)
            y = Inches(3.2)
        if sl.get("bullets"):
            text(s, Inches(0.7), y, Inches(12), Inches(4.5),
                 [f"•  {b}" for b in sl["bullets"]], 16, ink)
        tbl = sl.get("table")
        if tbl and tbl.get("headers"):
            rows = [tbl["headers"]] + tbl.get("rows", [])
            n_r, n_c = len(rows), len(tbl["headers"])
            gtbl = s.shapes.add_table(n_r, n_c, Inches(0.6), y, Inches(12.1),
                                      Inches(min(4.5, 0.45 * n_r))).table
            for ci in range(n_c):
                c = gtbl.cell(0, ci); c.text = str(tbl["headers"][ci])
                c.fill.solid(); c.fill.fore_color.rgb = primary
                c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                c.text_frame.paragraphs[0].font.size = Pt(11); c.text_frame.paragraphs[0].font.bold = True
            for ri, row in enumerate(tbl.get("rows", []), start=1):
                for ci in range(n_c):
                    val = row[ci] if ci < len(row) else ""
                    cell = gtbl.cell(ri, ci); cell.text = str(val)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return f"[created deck: {out.relative_to(pack_dir).as_posix()}]"


# ---------------------------------------------------------------------------
# Document (PDF) -- narrative docs incl. retraining plans
# ---------------------------------------------------------------------------

def build_doc(title: str, spec: dict, brand: dict, pack_dir: Path, out: Path) -> str:
    try:
        from fpdf import FPDF
    except ImportError:
        return "[document skipped: run `pip install fpdf2` to build PDFs]"

    primary = _rgb_tuple(brand["primary"])
    accent = _rgb_tuple(brand["accent"])
    company = brand.get("company", "")
    footer = brand.get("confidentiality_footer", "")
    classif = brand.get("classification", "")
    logo = _logo_path(pack_dir, brand)

    class Doc(FPDF):
        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "", 8)
            self.set_text_color(120, 120, 120)
            foot = " / ".join(p for p in (company, footer) if p)
            self.cell(0, 8, foot, align="L")
            self.cell(0, 8, f"Page {self.page_no()}", align="R")
        def header(self):
            if classif and self.page_no() > 1:
                self.set_font("Helvetica", "B", 8)
                self.set_text_color(*accent)
                self.cell(0, 6, classif.upper(), align="R", new_x="LMARGIN", new_y="NEXT")

    pdf = Doc()
    pdf.set_auto_page_break(True, margin=20)
    pdf.add_page()
    if logo:
        try: pdf.image(str(logo), x=15, y=12, h=14)
        except Exception: pass
    pdf.ln(28)
    if classif:
        pdf.set_font("Helvetica", "B", 9); pdf.set_text_color(*accent)
        pdf.cell(0, 6, classif.upper(), new_x="LMARGIN", new_y="NEXT")
    nl = dict(new_x="LMARGIN", new_y="NEXT")  # reset cursor to the left margin
    pdf.set_font("Helvetica", "B", 24); pdf.set_text_color(*primary)
    pdf.multi_cell(0, 11, title, **nl)
    if spec.get("subtitle"):
        pdf.set_font("Helvetica", "", 13); pdf.set_text_color(60, 60, 60)
        pdf.multi_cell(0, 7, spec["subtitle"], **nl)
    pdf.ln(4)

    for sec in spec.get("sections", []):
        pdf.set_font("Helvetica", "B", 14); pdf.set_text_color(*primary)
        pdf.multi_cell(0, 8, sec.get("heading", ""), **nl); pdf.ln(1)
        pdf.set_font("Helvetica", "", 11); pdf.set_text_color(30, 30, 30)
        if sec.get("body"):
            pdf.multi_cell(0, 6, sec["body"], **nl); pdf.ln(1)
        for b in sec.get("bullets", []):
            pdf.multi_cell(0, 6, f"  -  {b}", **nl)
        tbl = sec.get("table")
        if tbl and tbl.get("headers"):
            heads = tbl["headers"]; w = (pdf.w - 30) / max(1, len(heads))
            pdf.set_font("Helvetica", "B", 10); pdf.set_fill_color(*primary); pdf.set_text_color(255, 255, 255)
            for h in heads: pdf.cell(w, 7, str(h), border=0, fill=True)
            pdf.ln(7); pdf.set_text_color(30, 30, 30); pdf.set_font("Helvetica", "", 9)
            for row in tbl.get("rows", []):
                for ci in range(len(heads)):
                    pdf.cell(w, 6, str(row[ci]) if ci < len(row) else "", border="B")
                pdf.ln(6)
        pdf.ln(3)

    out.parent.mkdir(parents=True, exist_ok=True)
    pdf.output(str(out))
    return f"[created document: {out.relative_to(pack_dir).as_posix()}]"


# ---------------------------------------------------------------------------
# Financial model / projection (XLSX)
# ---------------------------------------------------------------------------

def build_sheet(title: str, spec: dict, brand: dict, pack_dir: Path, out: Path) -> str:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        return "[financial model skipped: run `pip install openpyxl` to build spreadsheets]"

    primary = _hex(brand["primary"])
    company = brand.get("company", "")
    classif = brand.get("classification", "")
    header_fill = PatternFill("solid", fgColor=primary)
    header_font = Font(bold=True, color="FFFFFF")

    wb = Workbook()
    first = True
    for sh in spec.get("sheets", []) or [{"name": "Sheet1", "headers": [], "rows": []}]:
        ws = wb.active if first else wb.create_sheet()
        first = False
        ws.title = (sh.get("name") or "Sheet")[:31]
        r = 1
        banner = " | ".join(p for p in (company, title, classif.upper() if classif else "") if p)
        if banner:
            ws.cell(r, 1, banner).font = Font(bold=True, size=12)
            r += 2
        headers = sh.get("headers", [])
        if headers:
            for ci, h in enumerate(headers, start=1):
                c = ws.cell(r, ci, str(h)); c.fill = header_fill; c.font = header_font
            r += 1
        for row in sh.get("rows", []):
            for ci, val in enumerate(row, start=1):
                ws.cell(r, ci, val)
            r += 1
        for ci in range(1, max(1, len(headers)) + 1):
            ws.column_dimensions[chr(64 + ci)].width = 22

    out.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out))
    return f"[created financial model: {out.relative_to(pack_dir).as_posix()}]"


# ---------------------------------------------------------------------------
# Autonomy dashboard (HTML) -- the standing scorecard for gate reviews
# ---------------------------------------------------------------------------
# Pure stdlib: the dashboard must build on any machine with zero optional
# libraries, because it is the artifact a buyer opens most often. Aletheia
# regenerates it whenever the numbers move; the file is self-contained and
# safe to email or drop into a shared drive.

def build_dashboard(title: str, spec: dict, brand: dict, pack_dir: Path, out: Path) -> str:
    import html as _html

    def esc(v) -> str:
        return _html.escape(str(v if v is not None else ""))

    primary = "#" + _hex(brand["primary"])
    accent = "#" + _hex(brand["accent"])
    company = brand.get("company", "")
    footer = brand.get("confidentiality_footer", "")
    classif = brand.get("classification", "")

    stage = spec.get("stage") or {}
    savings = spec.get("savings") or {}
    workflows = spec.get("workflows") or []
    rungs = ["L0", "L1", "L2", "L3", "L4", "L5"]

    def rung_ladder(current: str, target: str) -> str:
        cells = []
        for r in rungs:
            style = "background:#e5e7eb;color:#374151;"
            if r == str(current or "").upper():
                style = f"background:{primary};color:#fff;font-weight:700;"
            elif r == str(target or "").upper():
                style = f"background:{accent};color:#fff;font-weight:700;"
            cells.append(f'<td style="padding:4px 10px;text-align:center;{style}">{esc(r)}</td>')
        return "<table style='border-collapse:collapse'><tr>" + "".join(cells) + "</tr></table>"

    gate_color = {"passed": "#177245", "on_track": "#177245", "at_risk": "#B7791F",
                  "blocked": "#B03A2E", "pending": "#4B5563"}

    def li_block(heading: str, items: list) -> str:
        if not items:
            return ""
        lis = "".join(f"<li>{esc(i)}</li>" for i in items)
        return (f'<h2 style="color:{primary};font-size:16px;margin:24px 0 8px">{esc(heading)}</h2>'
                f'<ul style="margin:0;padding-left:20px;line-height:1.6">{lis}</ul>')

    wf_rows = []
    for w in workflows:
        status = str(w.get("gate_status") or "pending").lower()
        color = gate_color.get(status, "#4B5563")
        wf_rows.append(
            "<tr>"
            f'<td style="padding:6px 10px;border-bottom:1px solid #e5e7eb">{esc(w.get("name"))}</td>'
            f'<td style="padding:6px 10px;border-bottom:1px solid #e5e7eb;text-align:center">{esc(w.get("rung"))}'
            f' &rarr; {esc(w.get("target_rung"))}</td>'
            f'<td style="padding:6px 10px;border-bottom:1px solid #e5e7eb">{esc(w.get("gate"))}</td>'
            f'<td style="padding:6px 10px;border-bottom:1px solid #e5e7eb;color:{color};font-weight:600">'
            f'{esc(status.replace("_", " "))}</td>'
            f'<td style="padding:6px 10px;border-bottom:1px solid #e5e7eb">{esc(w.get("owner"))}</td>'
            "</tr>")
    wf_table = ""
    if wf_rows:
        wf_table = (
            f'<h2 style="color:{primary};font-size:16px;margin:24px 0 8px">Workflows</h2>'
            '<table style="border-collapse:collapse;width:100%;font-size:13px">'
            f'<tr style="background:{primary};color:#fff">'
            '<th style="padding:6px 10px;text-align:left">Workflow</th>'
            '<th style="padding:6px 10px">Rung</th>'
            '<th style="padding:6px 10px;text-align:left">Next gate</th>'
            '<th style="padding:6px 10px;text-align:left">Status</th>'
            '<th style="padding:6px 10px;text-align:left">Owner</th></tr>'
            + "".join(wf_rows) + "</table>")

    stat_cells = []
    for label, val in (
            ("Stage", f'{stage.get("current", "")} ({stage.get("score", "")})'),
            ("Target", f'{stage.get("target", "")} ({stage.get("target_score", "")}) by {stage.get("target_date", "")}'),
            ("Projected annual savings", savings.get("projected_annual", "")),
            ("Realized to date", savings.get("realized_to_date", ""))):
        stat_cells.append(
            '<div style="flex:1;min-width:160px;border:1px solid #e5e7eb;border-radius:8px;padding:12px 16px">'
            f'<div style="font-size:11px;text-transform:uppercase;letter-spacing:.05em;color:#6b7280">{esc(label)}</div>'
            f'<div style="font-size:18px;font-weight:700;color:{primary};margin-top:4px">{esc(val)}</div></div>')

    classif_line = (f'<div style="text-align:right;color:{accent};font-weight:700;'
                    f'font-size:11px;letter-spacing:.08em">{esc(classif.upper())}</div>') if classif else ""
    foot = " / ".join(p for p in (company, footer) if p)

    doc = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{esc(title)}</title></head>
<body style="font-family:{esc(brand.get('font') or 'Helvetica')},Arial,sans-serif;color:#1f2937;max-width:860px;margin:32px auto;padding:0 20px">
{classif_line}
<h1 style="color:{primary};font-size:26px;margin:8px 0 4px">{esc(title)}</h1>
<div style="color:#6b7280;font-size:13px;margin-bottom:20px">{esc(company)}</div>
<div style="display:flex;gap:12px;flex-wrap:wrap">{"".join(stat_cells)}</div>
<h2 style="color:{primary};font-size:16px;margin:24px 0 8px">Autonomy ladder</h2>
{rung_ladder(stage.get("current_rung") or stage.get("current"), stage.get("target_rung") or stage.get("target"))}
{wf_table}
{li_block("Wins", spec.get("wins") or [])}
{li_block("Risks", spec.get("risks") or [])}
{li_block("Next gates", spec.get("next_gates") or [])}
<div style="margin-top:36px;padding-top:12px;border-top:2px solid {accent};color:#6b7280;font-size:11px">{esc(foot)}</div>
</body></html>
"""
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(doc, encoding="utf-8")
    return f"[created dashboard: {out.relative_to(pack_dir).as_posix()}]"


BUILDERS = {"deck": build_deck, "doc": build_doc, "sheet": build_sheet,
            "dashboard": build_dashboard}
EXT = {"deck": "pptx", "doc": "pdf", "sheet": "xlsx", "dashboard": "html"}
